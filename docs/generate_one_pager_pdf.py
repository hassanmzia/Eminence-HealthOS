#!/usr/bin/env python3
"""Generate the Marketing 1-Pager as a styled PDF."""

import os

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
TEAL = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_BLUE = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x0F, 0x34, 0x60)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


def create_one_pager_pptx():
    """Create a single-slide one-pager as PPTX (can be exported to PDF)."""
    if not HAS_PPTX:
        print("python-pptx not installed, skipping PPTX one-pager")
        return

    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    def add_text(left, top, width, height, text, size=12, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return txBox

    def add_multiline(left, top, width, height, lines, size=9, color=DARK_GRAY, bold=False):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.space_after = Pt(3)
        return txBox

    def add_rect(left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # Header bar
    add_rect(0, 0, 8.5, 1.2, DARK_NAVY)
    add_text(0.5, 0.15, 7.5, 0.5, "EMINENCE HEALTHOS", 28, WHITE, True, PP_ALIGN.CENTER)
    add_text(0.5, 0.65, 7.5, 0.4, "The AI Operating System for Digital Healthcare Platforms", 14, TEAL, False, PP_ALIGN.CENTER)

    # Problem section
    y = 1.4
    add_text(0.5, y, 7.5, 0.3, "THE PROBLEM", 12, TEAL, True)
    add_text(0.5, y + 0.3, 7.5, 0.5,
             "Healthcare runs on 5+ disconnected systems costing $1.35T in annual waste. 86% of clinicians report burnout. No AI-native platform exists to unify clinical workflows.",
             9, DARK_GRAY)

    # Solution section
    y = 2.2
    add_text(0.5, y, 7.5, 0.3, "THE SOLUTION", 12, TEAL, True)
    add_text(0.5, y + 0.3, 7.5, 0.4,
             "HealthOS unifies RPM, Telehealth, Operations, and Analytics through 30 specialized AI agents in one platform.",
             10, DARK_GRAY, True)

    # Key Capabilities
    y = 2.9
    add_text(0.5, y, 7.5, 0.3, "KEY CAPABILITIES", 12, TEAL, True)

    caps = [
        "Remote Patient Monitoring — AI anomaly detection, trend analysis, smart alerts",
        "Telehealth — Video visits, ambient AI documentation, automated follow-up",
        "Operations Automation — Scheduling, prior auth, referrals, revenue cycle",
        "Population Health — Cohort analysis, readmission prediction, executive dashboards",
        "15 Total Modules — Digital Twin, Pharmacy, Labs, Imaging, Mental Health, and more",
    ]
    add_multiline(0.5, y + 0.3, 7.5, 1.2, caps, 9, DARK_GRAY)

    # Why HealthOS
    y = 4.4
    add_text(0.5, y, 7.5, 0.3, "WHY HEALTHOS?", 12, TEAL, True)
    whys = [
        "AI-Native — 30 specialized agents across 5 layers, not a chatbot add-on",
        "Unified — One platform replaces 5+ systems, 40-60% TCO reduction",
        "HIPAA by Design — Zero-trust, AES-256-GCM encryption, full audit trails",
        "Interoperable — FHIR R4 native, integrates with any EHR/pharmacy/lab",
        "Enterprise Scale — Kubernetes-native, 100 to 100,000+ patients",
    ]
    add_multiline(0.5, y + 0.3, 7.5, 1.2, whys, 9, DARK_GRAY)

    # Impact metrics
    y = 5.8
    add_rect(0.5, y, 7.5, 1.0, ACCENT_BLUE)
    add_text(0.5, y + 0.05, 7.5, 0.3, "MEASURABLE IMPACT", 11, WHITE, True, PP_ALIGN.CENTER)

    metrics = [
        ("60%+", "Admin\nAutomation"),
        ("30%", "Readmission\nReduction"),
        ("70%", "Less Doc\nTime"),
        ("40-60%", "Cost\nReduction"),
    ]
    for i, (val, label) in enumerate(metrics):
        x = 0.8 + i * 1.9
        add_text(x, y + 0.3, 1.5, 0.3, val, 18, WHITE, True, PP_ALIGN.CENTER)
        add_text(x, y + 0.6, 1.5, 0.4, label, 8, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # Technology
    y = 7.0
    add_text(0.5, y, 7.5, 0.3, "TECHNOLOGY", 12, TEAL, True)
    tech = [
        "Backend: Python 3.12, FastAPI, Temporal, Kafka | AI: Claude, GPT-4, Ollama, XGBoost, PyTorch",
        "Data: PostgreSQL, Redis, Qdrant, Neo4j | Frontend: React, Next.js 15, TypeScript",
        "Infra: Docker, Kubernetes, Helm, Terraform | CI/CD: GitHub Actions",
    ]
    add_multiline(0.5, y + 0.25, 7.5, 0.8, tech, 8, DARK_GRAY)

    # Business Model
    y = 8.0
    add_text(0.5, y, 3.5, 0.3, "BUSINESS MODEL", 12, TEAL, True)
    biz = [
        "Platform License (SaaS) — Annual subscription",
        "Module Add-Ons — Start RPM, expand to full",
        "Implementation Services — Custom integrations",
        "AI Marketplace — Third-party model revenue share",
    ]
    add_multiline(0.5, y + 0.3, 3.5, 1.0, biz, 8, DARK_GRAY)

    # Target Customers
    add_text(4.5, y, 3.5, 0.3, "TARGET CUSTOMERS", 12, TEAL, True)
    targets = [
        "Mid-size health systems (50-500 beds)",
        "Multi-site physician groups",
        "Accountable Care Organizations",
        "Digital health companies",
    ]
    add_multiline(4.5, y + 0.3, 3.5, 1.0, targets, 8, DARK_GRAY)

    # Footer
    add_rect(0, 9.8, 8.5, 1.2, DARK_NAVY)
    add_text(0.5, 9.95, 7.5, 0.3, "EMINENCE TECH SOLUTIONS", 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(0.5, 10.3, 7.5, 0.3, "sales@eminencetech.com | www.eminencetech.com", 10, TEAL, False, PP_ALIGN.CENTER)
    add_text(0.5, 10.6, 7.5, 0.3, "Intelligent Healthcare, Unified.", 9, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    output = os.path.join(os.path.dirname(__file__), "Eminence_HealthOS_One_Pager.pptx")
    prs.save(output)
    print(f"Created: {output}")


if __name__ == "__main__":
    create_one_pager_pptx()
