#!/usr/bin/env python3
"""Generate PowerPoint decks for Eminence HealthOS."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
MEDIUM_BLUE = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x0F, 0x34, 0x60)
TEAL = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_RED = RGBColor(0xE5, 0x39, 0x35)


def add_bg(slide, color=DARK_NAVY):
    """Add a solid background color to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, opacity=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=14, color=WHITE):
    """Add bulleted text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


# =============================================================================
# 1. PLATFORM OVERVIEW DECK
# =============================================================================
def create_platform_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Eminence HealthOS", 48, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                 "The AI Operating System for Digital Healthcare Platforms", 24, TEAL, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                 "30 Specialized AI Agents | 15 Clinical Modules | Enterprise-Grade Security",
                 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(6), Inches(11), Inches(0.5),
                 "Eminence Tech Solutions | 2026", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # --- Slide 2: The Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "The Problem: Fragmented Healthcare IT", 36, DARK_NAVY, True)

    problems = [
        "Separate RPM, telehealth, analytics, and billing systems that don't communicate",
        "Manual workflows causing delays in care coordination and prior authorizations",
        "Data silos preventing holistic patient views and population-level insights",
        "No AI-native architecture — AI bolted on as afterthought, not core design",
        "Compliance burden across multiple disconnected systems",
        "Clinician burnout from excessive documentation and administrative tasks",
    ]
    add_bullet_text(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(4.5), problems, 16, DARK_GRAY)

    # Stats box
    shape = add_shape(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(4.5), ACCENT_BLUE)
    add_text_box(slide, Inches(7.8), Inches(2.0), Inches(4.5), Inches(0.5),
                 "By the Numbers", 22, WHITE, True, PP_ALIGN.CENTER)
    stats = [
        "$4.5T — US healthcare spending annually",
        "30% — Administrative waste in healthcare",
        "86% — Clinicians report burnout symptoms",
        "2hrs — Documentation per 1hr of patient care",
        "$250B — Potential AI savings in healthcare",
    ]
    add_bullet_text(slide, Inches(7.8), Inches(2.8), Inches(4.5), Inches(3.5), stats, 15, WHITE)

    # --- Slide 3: The Solution ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "The Solution: Eminence HealthOS", 36, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.8),
                 "A unified AI operating system where 30+ specialized agents coordinate to deliver intelligent, automated healthcare workflows.",
                 18, TEAL)

    # Module cards
    modules = [
        ("Remote Patient\nMonitoring", "Continuous vitals\nAI anomaly detection\nSmart thresholds"),
        ("Telehealth\nPlatform", "Video consultations\nAmbient documentation\nAI visit prep"),
        ("Operations\nAutomation", "Smart scheduling\nPrior auth automation\nReferral coordination"),
        ("Population Health\nAnalytics", "Cohort segmentation\nReadmission prediction\nExecutive insights"),
    ]
    for i, (title, desc) in enumerate(modules):
        x = Inches(0.8 + i * 3.1)
        shape = add_shape(slide, x, Inches(2.8), Inches(2.8), Inches(3.8), MEDIUM_BLUE)
        add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(1),
                     title, 18, TEAL, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(2.4), Inches(2),
                     desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # --- Slide 4: 30-Agent Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "30-Agent Multi-Agent Architecture", 36, DARK_NAVY, True)

    layers = [
        ("Layer 5: Measurement", "Outcome Measurement | Population Health | Readmission Risk | Cohort Segmentation | Cost/Risk Insight | Executive Insight", RGBColor(0xF8, 0xCE, 0xCC)),
        ("Layer 4: Action", "Patient Communication | Scheduling | Prior Auth | Referral Coordination | Follow-Up Plan | Task Orchestration | Billing Readiness", RGBColor(0xFF, 0xF2, 0xCC)),
        ("Layer 3: Decisioning", "Master Orchestrator | Context Assembly | Policy/Rules Engine | Human-in-the-Loop | Audit/Trace | Quality/Confidence", RGBColor(0xD5, 0xE8, 0xD4)),
        ("Layer 2: Interpretation", "Anomaly Detection | Trend Analysis | Risk Scoring | Medication Review | Visit Preparation | Clinical Note Generation", RGBColor(0xDA, 0xE8, 0xFC)),
        ("Layer 1: Sensing", "Device Ingestion | Vitals Normalization | Adherence Monitoring | Insurance Verification | Escalation Routing", RGBColor(0xE1, 0xD5, 0xE7)),
    ]
    for i, (name, agents, color) in enumerate(layers):
        y = Inches(1.5 + i * 1.1)
        shape = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.0), color)
        add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3), Inches(0.4),
                     name, 16, DARK_NAVY, True)
        add_text_box(slide, Inches(1.0), y + Inches(0.45), Inches(11), Inches(0.5),
                     agents, 12, DARK_GRAY)

    # --- Slide 5: Technology Stack ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Technology Stack", 36, WHITE, True)

    stacks = [
        ("Backend", ["Python 3.12 + FastAPI", "Temporal Workflows", "Kafka Event Streaming", "Keycloak Auth"]),
        ("AI / ML", ["Claude API (Anthropic)", "OpenAI GPT-4", "Ollama (Local LLM)", "XGBoost, PyTorch"]),
        ("Data Layer", ["PostgreSQL + pgvector", "Redis 7 Cache", "Qdrant Vector DB", "Neo4j Knowledge Graph"]),
        ("Frontend", ["React 18 + Next.js 15", "TypeScript", "Tailwind + shadcn/ui", "Recharts / D3.js"]),
        ("DevOps", ["Docker + Kubernetes", "Helm 3 + Terraform", "GitHub Actions CI/CD", "Prometheus + Grafana"]),
    ]
    for i, (title, items) in enumerate(stacks):
        x = Inches(0.5 + i * 2.5)
        add_shape(slide, x, Inches(1.6), Inches(2.3), Inches(4.5), MEDIUM_BLUE)
        add_text_box(slide, x + Inches(0.1), Inches(1.7), Inches(2.1), Inches(0.5),
                     title, 18, TEAL, True, PP_ALIGN.CENTER)
        add_bullet_text(slide, x + Inches(0.2), Inches(2.4), Inches(2.0), Inches(3.5), items, 13, LIGHT_GRAY)

    # --- Slide 6: All 15 Modules ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "15 Clinical & Operational Modules", 36, DARK_NAVY, True)

    all_modules = [
        "Remote Patient Monitoring", "Telehealth", "Operations", "Analytics",
        "Digital Twin", "Pharmacy", "Labs", "Medical Imaging",
        "Revenue Cycle (RCM)", "Mental Health", "Patient Engagement", "Compliance",
        "Ambient AI", "Research & Genomics", "AI Marketplace"
    ]
    colors = [TEAL, ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, RGBColor(0x9C, 0x27, 0xB0)] * 3
    for i, mod in enumerate(all_modules):
        row, col = divmod(i, 5)
        x = Inches(0.6 + col * 2.5)
        y = Inches(1.5 + row * 1.8)
        shape = add_shape(slide, x, y, Inches(2.3), Inches(1.4), colors[i % len(colors)])
        add_text_box(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.1), Inches(0.8),
                     mod, 15, WHITE, True, PP_ALIGN.CENTER)

    # --- Slide 7: Security & Compliance ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Enterprise Security & Compliance", 36, WHITE, True)

    sec_items = [
        ("Zero Trust Architecture", "Every request verified, no implicit trust"),
        ("HIPAA Compliance", "PHI encryption (AES-256-GCM), audit trails, access controls"),
        ("RBAC Authorization", "22 granular permissions across 6 clinical roles"),
        ("Multi-Tenant Isolation", "Data, query, agent, and API-level isolation"),
        ("PHI Protection", "5-level classification, automated detection & masking"),
        ("Full Audit Trail", "Complete decision chain traceability for all AI actions"),
        ("Input Sanitization", "SQL injection, XSS, and command injection prevention"),
        ("Regulatory Ready", "HITRUST, SOC2, FDA SaMD, EU AI Act compliance paths"),
    ]
    for i, (title, desc) in enumerate(sec_items):
        row, col = divmod(i, 2)
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.5 + row * 1.4)
        add_shape(slide, x, y, Inches(5.8), Inches(1.2), MEDIUM_BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.4),
                     title, 16, TEAL, True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.4), Inches(0.5),
                     desc, 13, LIGHT_GRAY)

    # --- Slide 8: Deployment ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Deployment & Infrastructure", 36, DARK_NAVY, True)

    deploy_items = [
        "Docker Compose for local development with hot-reload",
        "Kubernetes (Helm 3) for staging and production",
        "Horizontal Pod Autoscaler: 3-10 pods, 60% CPU / 70% memory targets",
        "CI/CD via GitHub Actions: lint, test, type-check, security scan",
        "Prometheus + Grafana monitoring with custom dashboards",
        "OpenTelemetry distributed tracing across all agents",
        "Temporal workflow engine for long-running healthcare processes",
        "Automated deploy/teardown scripts with rollback support",
    ]
    add_bullet_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(5), deploy_items, 18, DARK_GRAY)

    # --- Slide 9: Roadmap ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "5-Year Platform Roadmap", 36, WHITE, True)

    roadmap = [
        ("Year 1", "Remote Patient Monitoring MVP", "Core agent framework, RPM module, vitals pipeline, clinician dashboard"),
        ("Year 2", "Telehealth Platform", "Video consultations, ambient AI, clinical note generation"),
        ("Year 3", "Operations Automation", "Scheduling, prior auth, referrals, billing workflows"),
        ("Year 4", "Population Health Analytics", "Cohort analysis, readmission prediction, executive dashboards"),
        ("Year 5", "Autonomous Operations", "Full platform integration, AI marketplace, research & genomics"),
    ]
    for i, (year, phase, desc) in enumerate(roadmap):
        y = Inches(1.5 + i * 1.15)
        add_shape(slide, Inches(0.8), y, Inches(1.5), Inches(0.9), TEAL)
        add_text_box(slide, Inches(0.9), y + Inches(0.15), Inches(1.3), Inches(0.6),
                     year, 20, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(2.6), y + Inches(0.05), Inches(4), Inches(0.4),
                     phase, 18, TEAL, True)
        add_text_box(slide, Inches(2.6), y + Inches(0.45), Inches(9), Inches(0.4),
                     desc, 13, LIGHT_GRAY)

    # --- Slide 10: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "Eminence HealthOS", 48, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 "The AI Operating System for Digital Healthcare Platforms", 24, TEAL, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(1),
                 "Eminence Tech Solutions\ncontact@eminencetech.com", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    output = os.path.join(os.path.dirname(__file__), "Eminence_HealthOS_Platform_Overview.pptx")
    prs.save(output)
    print(f"Created: {output}")


# =============================================================================
# 2. MARKETING PITCH DECK
# =============================================================================
def create_pitch_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(1), Inches(1), Inches(11), Inches(1.5),
                 "Eminence HealthOS", 52, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
                 "The AI Operating System for Digital Healthcare", 28, TEAL, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                 "Investor Pitch Deck | 2026", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "CONFIDENTIAL", 14, ACCENT_RED, True, PP_ALIGN.CENTER)

    # --- Slide 2: Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(1),
                 "The $4.5T Problem", 36, DARK_NAVY, True)
    problems = [
        "Healthcare IT is a $350B+ market of disconnected point solutions",
        "30% of healthcare spending is administrative waste ($1.35T)",
        "Clinicians spend 2 hours on documentation for every 1 hour with patients",
        "No unified AI platform exists to orchestrate clinical workflows end-to-end",
        "Existing solutions bolt on AI — they weren't designed AI-first",
    ]
    add_bullet_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(4.5), problems, 18, DARK_GRAY)

    # --- Slide 3: Solution ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Our Solution", 36, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.8),
                 "HealthOS is the first AI-native operating system that unifies RPM, Telehealth, Operations, and Analytics through 30 specialized AI agents.",
                 20, TEAL)

    points = [
        ("Unified Platform", "One system replaces 5+ disconnected tools"),
        ("30 AI Agents", "Specialized agents for every clinical workflow"),
        ("HIPAA Native", "Zero-trust security, PHI encryption, full audit trails"),
        ("Enterprise Ready", "Multi-tenant, Kubernetes-native, horizontally scalable"),
    ]
    for i, (title, desc) in enumerate(points):
        x = Inches(0.5 + i * 3.15)
        add_shape(slide, x, Inches(3), Inches(2.9), Inches(3.5), MEDIUM_BLUE)
        add_text_box(slide, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(0.5),
                     title, 20, TEAL, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(2.5), Inches(2),
                     desc, 15, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # --- Slide 4: Market Opportunity ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Market Opportunity", 36, DARK_NAVY, True)

    markets = [
        ("TAM", "$350B", "Global Healthcare IT Market"),
        ("SAM", "$45B", "AI in Healthcare Market (2026)"),
        ("SOM", "$2B", "AI-Native Healthcare Platforms"),
    ]
    for i, (label, value, desc) in enumerate(markets):
        x = Inches(1 + i * 4)
        add_shape(slide, x, Inches(1.8), Inches(3.5), Inches(3), ACCENT_BLUE)
        add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(3.1), Inches(0.6),
                     label, 24, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(3.1), Inches(1),
                     value, 48, TEAL, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(3.8), Inches(3.1), Inches(0.8),
                     desc, 14, WHITE, False, PP_ALIGN.CENTER)

    growth = [
        "Healthcare AI market growing at 45% CAGR",
        "RPM market alone projected at $175B by 2030",
        "Telehealth adoption permanently elevated post-pandemic",
        "Regulatory tailwinds: CMS RPM reimbursement expansion, ONC interoperability mandates",
    ]
    add_bullet_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(2), growth, 15, DARK_GRAY)

    # --- Slide 5: Product ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Product: 15 Modules, One Platform", 36, WHITE, True)

    product_cols = [
        ("Core Clinical", ["Remote Patient Monitoring", "Telehealth & Video", "Digital Twin", "Ambient AI Documentation"]),
        ("Operations", ["Smart Scheduling", "Prior Authorization", "Revenue Cycle Mgmt", "Referral Coordination"]),
        ("Intelligence", ["Population Health Analytics", "Risk Prediction", "Cohort Segmentation", "Executive Insights"]),
        ("Specialty", ["Pharmacy & Medications", "Labs Integration", "Medical Imaging", "Mental Health"]),
    ]
    for i, (title, items) in enumerate(product_cols):
        x = Inches(0.5 + i * 3.15)
        add_shape(slide, x, Inches(1.6), Inches(2.9), Inches(5), MEDIUM_BLUE)
        add_text_box(slide, x + Inches(0.1), Inches(1.7), Inches(2.7), Inches(0.5),
                     title, 18, TEAL, True, PP_ALIGN.CENTER)
        add_bullet_text(slide, x + Inches(0.2), Inches(2.4), Inches(2.6), Inches(4), items, 14, LIGHT_GRAY)

    # --- Slide 6: Competitive Advantage ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Why HealthOS Wins", 36, DARK_NAVY, True)

    advantages = [
        ("AI-Native Architecture", "30 specialized agents vs. generic AI add-ons. Purpose-built multi-agent system is the core technical moat."),
        ("Unified Platform", "One platform replaces 5+ point solutions. Lower TCO, unified data, consistent experience."),
        ("Enterprise Security", "Zero-trust, HIPAA-native, multi-tenant. Not bolted-on compliance — built from day one."),
        ("Interoperability", "FHIR R4 native with HL7, NCPDP, X12 support. Seamless EHR integration."),
        ("Scalable Architecture", "Kubernetes-native, event-driven, horizontally scalable. From 100 to 100,000+ patients."),
        ("Speed to Market", "Modular architecture enables rapid deployment. New modules in weeks, not months."),
    ]
    for i, (title, desc) in enumerate(advantages):
        row, col = divmod(i, 2)
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.5 + row * 1.8)
        add_shape(slide, x, y, Inches(5.8), Inches(1.5), ACCENT_BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.5),
                     title, 18, WHITE, True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(5.4), Inches(0.8),
                     desc, 13, WHITE)

    # --- Slide 7: Business Model ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Business Model", 36, WHITE, True)

    biz_items = [
        ("Platform License (SaaS)", "Annual subscription per organization. Tiered by modules and patient volume."),
        ("Module Add-Ons", "Additional modules purchased separately. Start with RPM, expand to full platform."),
        ("Implementation Services", "Custom EHR integrations, data migration, workflow configuration. Billable professional services."),
        ("AI Marketplace Revenue", "Revenue share from third-party AI models deployed on the platform."),
    ]
    for i, (title, desc) in enumerate(biz_items):
        y = Inches(1.5 + i * 1.4)
        add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.2), MEDIUM_BLUE)
        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(4), Inches(0.4),
                     title, 20, TEAL, True)
        add_text_box(slide, Inches(1.2), y + Inches(0.55), Inches(10.5), Inches(0.5),
                     desc, 14, LIGHT_GRAY)

    # --- Slide 8: Go-to-Market ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Go-to-Market Strategy", 36, DARK_NAVY, True)

    gtm = [
        ("Phase 1: Land", "Target mid-size health systems (50-500 beds) with RPM module as entry point. Value prop: reduce readmissions by 30%, automate 60% of vitals monitoring."),
        ("Phase 2: Expand", "Upsell telehealth + operations modules to existing RPM customers. Cross-sell analytics for population health."),
        ("Phase 3: Platform", "Full platform deployment. AI marketplace for third-party models. Research & genomics modules for academic medical centers."),
    ]
    for i, (title, desc) in enumerate(gtm):
        y = Inches(1.5 + i * 1.8)
        add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.5), ACCENT_BLUE)
        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.4),
                     title, 20, WHITE, True)
        add_text_box(slide, Inches(1.2), y + Inches(0.55), Inches(10.5), Inches(0.8),
                     desc, 14, WHITE)

    # --- Slide 9: Team / Company ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "About Eminence Tech Solutions", 36, WHITE, True)
    company_points = [
        "AI Platform Vendor — not a contractor or consultant",
        "Full ownership of core platform IP (agent orchestration, ML pipelines, security framework)",
        "Licensable enterprise SaaS with protected intellectual property",
        "Three-layer product boundary: Core Platform (IP) | Licensed Modules | Client Services",
        "Built by engineers with deep expertise in healthcare AI, distributed systems, and clinical workflows",
    ]
    add_bullet_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(4.5), company_points, 18, LIGHT_GRAY)

    # --- Slide 10: Ask ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Let's Build the Future\nof Healthcare AI Together", 44, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4), Inches(11), Inches(1),
                 "Eminence HealthOS — The AI Operating System for Digital Healthcare Platforms", 20, TEAL, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
                 "contact@eminencetech.com | www.eminencetech.com", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    output = os.path.join(os.path.dirname(__file__), "Eminence_HealthOS_Pitch_Deck.pptx")
    prs.save(output)
    print(f"Created: {output}")


# =============================================================================
# 3. MARKETING DECK
# =============================================================================
def create_marketing_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Eminence HealthOS", 52, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
                 "Intelligent Healthcare, Unified", 28, TEAL, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                 "Product Marketing Deck | 2026", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # --- Slide 2: What is HealthOS ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "What is Eminence HealthOS?", 36, DARK_NAVY, True)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
                 "HealthOS is the first AI-native healthcare operating system. It unifies remote patient monitoring, telehealth, operations automation, and population health analytics into a single intelligent platform powered by 30 specialized AI agents.",
                 20, DARK_GRAY)

    value_props = [
        ("Reduce Costs", "Automate 60%+ of administrative tasks. Reduce readmissions by up to 30%."),
        ("Improve Outcomes", "AI-driven anomaly detection catches deterioration days earlier."),
        ("Unify Workflows", "One platform replaces 5+ disconnected systems."),
    ]
    for i, (title, desc) in enumerate(value_props):
        x = Inches(0.8 + i * 4)
        add_shape(slide, x, Inches(3.8), Inches(3.7), Inches(2.8), ACCENT_BLUE)
        add_text_box(slide, x + Inches(0.2), Inches(4.0), Inches(3.3), Inches(0.5),
                     title, 22, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(4.8), Inches(3.3), Inches(1.5),
                     desc, 15, WHITE, False, PP_ALIGN.CENTER)

    # --- Slide 3: For Clinicians ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "For Clinicians & Care Teams", 36, WHITE, True)
    clinician_features = [
        "AI-prepared visit summaries before every consultation",
        "Real-time vitals monitoring with intelligent anomaly alerts",
        "Ambient AI documentation — notes write themselves",
        "One-click lab orders, prescriptions, and referrals",
        "Risk-stratified patient panels with actionable insights",
        "Telehealth with integrated clinical decision support",
        "Automated follow-up plans and care coordination",
    ]
    add_bullet_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), clinician_features, 18, LIGHT_GRAY)

    # --- Slide 4: For Health Systems ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "For Health Systems & Administrators", 36, DARK_NAVY, True)
    admin_features = [
        "Population health analytics with executive dashboards",
        "Automated prior authorizations and insurance verification",
        "Revenue cycle optimization with AI-powered coding",
        "Multi-tenant architecture for health system networks",
        "HIPAA-compliant with full audit trail and compliance reporting",
        "Scalable from 100 to 100,000+ patients",
        "Reduce administrative FTEs by 40-60%",
    ]
    add_bullet_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), admin_features, 18, DARK_GRAY)

    # --- Slide 5: Key Differentiators ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Why HealthOS?", 36, WHITE, True)

    diffs = [
        ("30 AI Agents", "Not a single chatbot. 30 specialized agents working together across 5 operational layers."),
        ("AI-Native", "Built AI-first from the ground up. Not an EHR with AI bolted on."),
        ("One Platform", "RPM + Telehealth + Operations + Analytics in one unified system."),
        ("HIPAA by Design", "Zero-trust architecture, field-level PHI encryption, complete audit trails."),
        ("Interoperable", "FHIR R4 native. Integrates with any EHR, pharmacy, lab, or payer system."),
        ("Enterprise Scale", "Kubernetes-native, event-driven. Scales from clinic to health system."),
    ]
    for i, (title, desc) in enumerate(diffs):
        row, col = divmod(i, 3)
        x = Inches(0.5 + col * 4.15)
        y = Inches(1.5 + row * 2.7)
        add_shape(slide, x, y, Inches(3.9), Inches(2.3), MEDIUM_BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.5),
                     title, 20, TEAL, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.5), Inches(1.3),
                     desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # --- Slide 6: Getting Started ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                 "Getting Started", 36, DARK_NAVY, True)
    steps = [
        ("1. Discovery Call", "We understand your organization's needs, workflows, and integration requirements."),
        ("2. Platform Demo", "See HealthOS in action with your clinical scenarios."),
        ("3. Pilot Deployment", "Deploy a focused pilot (e.g., RPM for a specific patient cohort)."),
        ("4. Full Rollout", "Expand to full platform with additional modules and integrations."),
    ]
    for i, (title, desc) in enumerate(steps):
        y = Inches(1.5 + i * 1.4)
        add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.2), ACCENT_BLUE)
        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.4),
                     title, 20, WHITE, True)
        add_text_box(slide, Inches(1.2), y + Inches(0.55), Inches(10.5), Inches(0.5),
                     desc, 15, WHITE)

    # --- Slide 7: Contact ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "Ready to Transform\nYour Healthcare Operations?", 40, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4), Inches(11), Inches(1),
                 "Schedule a demo today", 24, TEAL, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
                 "sales@eminencetech.com | www.eminencetech.com", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    output = os.path.join(os.path.dirname(__file__), "Eminence_HealthOS_Marketing_Deck.pptx")
    prs.save(output)
    print(f"Created: {output}")


if __name__ == "__main__":
    create_platform_deck()
    create_pitch_deck()
    create_marketing_deck()
    print("\nAll decks generated successfully!")
